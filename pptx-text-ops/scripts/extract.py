"""Extract all text from a PPTX file for inspection and validation.

Prints slide text and speaker notes to stdout, one paragraph per line,
so the output can be piped to grep, awk, or other tools.

Usage:
    python extract.py <file.pptx>

Examples:
    python extract.py presentation.pptx
    python extract.py edited.pptx | grep -iE "xxxx|lorem|ipsum"
    python extract.py edited.pptx | awk 'length > 120'
"""

import argparse
import sys
from pathlib import Path
from pptx import Presentation


def extract(input_file: str) -> str:
    src = Path(input_file)
    if not src.exists():
        return f"Error: {input_file} does not exist"
    if src.suffix.lower() != ".pptx":
        return f"Error: {input_file} must be a .pptx file"

    try:
        prs = Presentation(str(src))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"<!-- slide {i} -->")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[notes] {notes}")
        return "\n".join(lines)
    except Exception as e:
        return f"Error: {e}"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text from a PPTX file")
    parser.add_argument("input_file", help="PPTX file to read")
    args = parser.parse_args()

    result = extract(args.input_file)
    print(result)
    if result.startswith("Error"):
        sys.exit(1)
